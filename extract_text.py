import sys
import io
from pptx import Presentation

# Reconfigure stdout to use UTF-8 encoding
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

def extract_text_from_pptx(file_path):
    """Extracts all text from a .pptx file."""
    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return "\n".join(text_runs)
    except Exception as e:
        # Return a string representation of the error
        return f"Error processing PPTX file: {e}"

if __name__ == "__main__":
    if len(sys.argv) > 1:
        filepath = sys.argv[1]
        # Print the result to stdout
        print(extract_text_from_pptx(filepath))
    else:
        print("Error: Please provide a file path as an argument.") 